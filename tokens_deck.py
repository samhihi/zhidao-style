# -*- coding: utf-8 -*-
"""知到 ZhiDao · 设计系统 · 演讲载体常量库+骨架函数(python-pptx)

知到 DS v1.1 · 2026-07-14 · 修订记录见 design-system.md

事实来源: 《从0到1，用AI Agent搭一个能跑的产品》分享 deck 41 页(私有制品, 未随本仓库分发)。
全部数值经逐页 XML 溯源校验(2026-07-14); 总纲见 design-system.md。
用法: 生成新 deck 时 `from tokens_deck import *`, 几何一律用这里的 EMU 常量,
文字一律走 set_font() 把 latin/ea/cs 三槽同时设成思源黑体;
开工先跑 assert_fonts_installed()——缺字体显式报错, 绝不静默落雅黑。
骨架件(压暗渐变纱/页眉三件套/结论条)一律调文件尾的 add_* 函数, 别逐页手放。
"""

from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

# ============ 画布(16:9, 13.33 x 7.5 in) ============
SLIDE_W = Emu(12192000)          # v1.1 定版: 真 16:9 整数值, 新 deck 一律用它
SLIDE_H = Emu(6858000)
SLIDE_W_LEGACY = Emu(12191365)   # 原片实测值, 系模板转换残留(presentation.xml 还标着
                                 # type="screen4x3"), 比真 16:9 窄 635 EMU——仅改旧片时用, 勿仿

# ============ 色板 · 冷深面 ============
BG_ABYSS   = RGBColor(0x03, 0x07, 0x0E)   # 渐变压暗最深档
BG_DEEP    = RGBColor(0x06, 0x11, 0x1A)   # 页底
BG_SHADE   = RGBColor(0x0C, 0x14, 0x18)   # 过渡档; 也是绿底表头的深色字
PANEL      = RGBColor(0x14, 0x20, 0x28)   # 面板/终端窗体/章封面高亮框底/结论条底
PANEL_2    = RGBColor(0x1B, 0x27, 0x35)   # 次面板/终端标题栏
BORDER     = RGBColor(0x24, 0x34, 0x41)   # 中性描边(1pt)

TEXT       = RGBColor(0xE6, 0xED, 0xF3)   # 主文字
TEXT_DIM   = RGBColor(0x8A, 0x97, 0xA3)   # 次文字(灰阶只有这两档, 别自创中间档)

BRAND      = RGBColor(0x16, 0xC7, 0x9A)   # 青绿: 骨架/正向/强调
GOLD       = RGBColor(0xE0, 0xB8, 0x4A)   # 暖金: 结论/金句(每页金色焦点尽量唯一)
GOLD_HOT   = RGBColor(0xFF, 0xC0, 0x00)   # 亮金: 全片只留给「临门一脚」(原片仅 2 页)

DOT_RED    = RGBColor(0xFF, 0x5F, 0x56)   # 终端窗三点(全片唯一的红也在这)
DOT_YELLOW = RGBColor(0xFF, 0xBD, 0x2E)
DOT_GREEN  = RGBColor(0x27, 0xC9, 0x3F)
# 禁用 #00B050 —— 原片两处属色板漂移, 一律用 BRAND

# ============ 字体 ============
FONT = "思源黑体"        # = Noto Sans SC / Source Han Sans, 与 Web 载体同源
# v1.1: 原 FONT_FALLBACK="微软雅黑" 死常量已删——缺字体不做静默兜底,
# 一律 assert_fonts_installed() 显式报错(堵死 typography 分册 Don't#4「混入雅黑」的复发路径)。

def set_font(run, size_pt, color=TEXT, bold=False, font=FONT):
    """latin/ea/cs 三槽一起设, 否则中英文/符号会各走各的字体。

    v1.1 修复: ea/cs 不再盲目 append——rPr 子元素有顺序约束(CT_TextCharacterProperties),
    必须插到 a:sym/a:hlinkClick/a:hlinkMouseOver/a:rtl/a:extLst 之前,
    否则带超链接的 run 会产出 schema 非法文件(PowerPoint 弹修复)。
    """
    f = run.font
    f.size, f.bold = Pt(size_pt), bold
    f.color.rgb = color
    f.name = font                       # latin 槽(python-pptx 自会插对位置)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.insert_element_before(
                e, "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")
        e.set("typeface", font)

# ============ 字号刻度(pt) · 原片 11.5-44 共 27 级中的骨干档 ============
SZ_THANKS      = 44     # 谢幕「谢谢大家」(全片最大)
SZ_COVER       = 40.5   # 封面主标题
SZ_DEMO        = 34     # 「▶ 现场演示」(正文区最大字)
SZ_CHAPTER     = 30     # 章封面章名
SZ_TITLE       = 29     # 内容页统一页题(29pt 档 33/41 页, 别动)
SZ_QUOTE_BOX   = 27     # 框体金句(S36=27/S40=28, 收口取 27; 恒低于页题)  v1.1 拆分
SZ_QUOTE_BARE  = 26     # 裸排金句(S5/S28)                              v1.1 拆分
SZ_QUOTE       = SZ_QUOTE_BARE  # @deprecated v1.0 旧名只为兼容保留, 新代码用上面两个
SZ_STATEMENT   = 24     # 大字陈述层 24-30
SZ_LEAD        = 19     # 引导行/终端窗内重点行
SZ_AGENDA      = 16     # 章封面「接下来」/议程正文/结论条引导语
SZ_BODY        = 15     # 卡片正文 14.5-15/结论条正文
SZ_PROGRESS    = 14     # 「第 N 章 · 共六章」(绿粗)
SZ_BREAD_CH    = 12.5   # 面包屑·章名(绿粗); 章封面小注同档
SZ_BREAD_SEC   = 11.5   # 面包屑·节名(灰), 全片最小字号——仅限装饰性路标, 观众无需读到
# 字距: 全片 0 处 spc 属性, 宽字距一律靠手工空格(如「第 一 章 · 共六章」)
# 行距: 表格 1.10 / 终端窗 1.25 / 长文 1.50(线间距用 line_spacing 设)

# ============ 页眉(骨架口径 v1.1 更正) ============
# 口径: 绿竖条 39/41(缺 S1/S41); 29pt 页题 33/41; 完整三件套约 32 页=内容页体例;
# 章封面用变体(加高竖条+30pt 章名+「第 N 章 · 共六章」进度语, 不挂面包屑);
# S1/S2/S41 破格, 别给破格页硬挂面包屑。
# 内容页一律调 add_header()(见文件尾)——原片 41 页逐页手放是最大技术债, 母版化留 v2。
HEADER_BAR_X, HEADER_BAR_Y = Emu(548640), Emu(566928)   # 左侧品牌绿竖条
HEADER_BAR_W, HEADER_BAR_H = Emu(100584), Emu(566928)   # 0.11 x 0.62 in, 实色无描边
HEADER_BAR_Y_CHAPTER = Emu(502920)                      # 章封面竖条起点上移档(slide3 实测)
HEADER_BAR_H_CHAPTER = Emu(603504)                      # 章封面竖条加高档
TITLE_X = Emu(777240)   # 页题框左缘(29pt 粗白); 框默认左内距 91440 后文字起于版心 868680
TITLE_Y = Emu(475488)   # 页题框 y·双行档(33 个 29pt 页题中 23 页的主档; 原框 spAutoFit,
                        # 单行页 y/cy 会浮动——真正恒定的是下面的垂直中心)
TITLE_W = Emu(10607040)                                 # slide18 实测 ext
TITLE_H = Emu(868680)                                   # 双行档框高(slide18 实测)
TITLE_CY_CENTER = Emu(909828)   # 页题垂直中心 = TITLE_Y + TITLE_H/2, 全片恒定锚点
BREAD_X, BREAD_Y, BREAD_W = Emu(4572000), Emu(201168), Emu(7178040)
BREAD_H = Emu(191770)                                   # 面包屑框高(slide18 实测, 单行)
# 面包屑右对齐: 「章名(12.5pt 绿粗)␣␣·␣␣节名(11.5pt 灰)」分隔=两个半角空格夹「·」
BREAD_SEP = "  ·  "

# ============ 组件规格(经三页/多页同构验证) ============
# 终端窗(演示页 S20/25/30 三页逐值相同)
TERM_X, TERM_Y = Emu(2441448), Emu(2331720)
TERM_W, TERM_H = Emu(7315200), Emu(3017520)      # 8.0 x 3.3 in
TERM_RADIUS_ADJ = 4000                            # roundRect adj
TERM_TITLEBAR_H = Emu(384048)                     # 标题栏 PANEL_2, 窗体 PANEL
TERM_BORDER_W = Pt(1)                             # BORDER 色
TERM_DOT_D, TERM_DOT_GAP = Emu(137160), Emu(237744)   # 三点 0.15in, 间距 0.26in
TERM_DOTS_X, TERM_DOTS_Y = Emu(2697480), Emu(2459736) # 三点起点(slide20 实测)

# 章封面当前章高亮框(六张章封面同构)
CHAP_BOX_X, CHAP_BOX_Y = Emu(658368), Emu(1627632)    # slide3 实测(六章同位)
CHAP_BOX_W, CHAP_BOX_H = Emu(5989320), Emu(594360)    # 6.55 x 0.65 in
CHAP_BOX_ADJ = 9000
CHAP_BOX_BORDER = Emu(15240)    # v1.1 定版 1.2pt: 品牌绿描边一律 1.2pt(描边收口, 见下);
                                # 原片实测 16510=1.3pt, 属 1.2/1.3 漂移中的多数档, 历史事实记档于此
# 框内: 「▶」19pt 绿粗 + 章名 18pt(当前章粗白, 其余 18pt 常规 TEXT_DIM)

# 金句框(S36/S40 同构, 宽度可参数化 70-76% 版面宽; PANEL 底 + 1.4pt GOLD 描边)
QUOTE_BORDER = Emu(17780)                                  # 1.4pt GOLD 描边
QUOTE_ADJ = 5000
QUOTE_X_WIDE, QUOTE_Y_WIDE = Emu(1463040), Emu(2926080)    # S36 实测(slide36.xml roundRect xfrm)
QUOTE_W_WIDE, QUOTE_H_WIDE = Emu(9235440), Emu(1737360)    # S36 档
QUOTE_X_SLIM, QUOTE_Y_SLIM = Emu(1828800), Emu(2697480)    # S40 实测(slide40.xml roundRect xfrm)
QUOTE_W_SLIM, QUOTE_H_SLIM = Emu(8503920), Emu(1097280)    # S40 档

# 原生表格(5 页同构: 左缘对齐版心, 行高 0.70-0.90in)
TABLE_X = Emu(868680)
TABLE_Y = Emu(1783080)   # 主档(slide14/slide18 <a:tbl> graphicFrame xfrm 同值);
                         # 上方有引导内容时下移: S32=1828800 / S24=1920240 / S12=2148840
TABLE_W = Emu(10469880)  # 5 页表格 graphicFrame 全部同宽(11.45in, 与结论条同宽)
# 表头: BRAND 底 + BG_SHADE 深色字 15pt 粗(双色表头则绿+金)

# 底部结论条(slide38 实测; y/高度随页内容可浮动, 此处为默认档)
CONCL_X, CONCL_Y = Emu(868680), Emu(5074920)
CONCL_W, CONCL_H = Emu(10469880), Emu(1051560)
CONCL_TEXT_INSET = Emu(274320)   # 条内文字左缩进 0.3in(S38: 文本框 1143000 - 条 868680)

# 谢幕页二维码容器
QR_BOX = Emu(2766060)                                 # 正方形, PANEL 底 + 1pt BORDER
QR_IMG = Emu(2407920)                                 # 四周留白约 0.20 in

# 内容底部安全线
BOTTOM_SAFE = Emu(6309360)   # 6.9in: 内容形状底缘别越过这条线, 再往下只留页脚气口

# ============ 描边(v1.1 收口) ============
# 语义主通道是颜色(中性 BORDER / 品牌绿 BRAND / 金 GOLD), 粗细只是制作规格三档——
# 0.1-0.2pt 差在投影距离不可感知, 别把层级压在粗细上。
# 三档: 中性 1.0pt / 品牌绿一律 1.2pt / 金句 1.4pt。
# 原片品牌绿实测漂移: 1.2pt x3 页(S18/19/28)、1.3pt x6 页(S3/9/16/22/27/34)、
# S38 1.0pt 属例外——收口取 1.2, 历史测量保留为事实。
STROKE_NEUTRAL, STROKE_BRAND, STROKE_GOLD = Pt(1), Pt(1.2), Pt(1.4)

# ============ 全幅底图 + 压暗渐变纱(背景三明治) ============
# 底图: AI 插画 2048x1152 为主(约 154dpi), 与画布同宽; 4K 放映升 3840x2160
# 纱: 一律调 add_veil()(v1.1 新增, 见下)——python-pptx 公开 API 只能做 2 停且不带停靠点
# 透明度, 含 alpha 的三停纱必须注入 <a:gradFill> OOXML, 已在函数里封装好。
# 恒定的只有两件事: ①只用 BG_ABYSS/BG_DEEP 两个深色 ②文字所在区域必须压到可读。
GRADIENT_ANGLE_DOWN = 5400000     # 1/60000 度; 90° 主流(38/41 页)
GRADIENT_ANGLE_UP   = 16200000    # 270° 反向(S7/S10/S20)

# 默认纱 = 后半场模态配置(S23/27/30/31/33/34/35 七页逐值一致):
# (色, 位置%, 不透明度%)——顶全遮 / 中透出插画 / 底压暗托文字; 中停 α 随插画亮度 5-59% 调
VEIL_STOPS_DEFAULT = [(BG_ABYSS, 0, 100), (BG_DEEP, 69, 25), (BG_ABYSS, 100, 84)]

# ============ 骨架函数(v1.1 新增: 单点维护, 别逐页手放) ============

def add_veil(slide, stops=None, ang=None):
    """压暗渐变纱(三层三明治第二层): 自建全幅无边框矩形 + 含 alpha 的多停渐变。

    stops: [(RGBColor, 位置%, 不透明度%), ...], 缺省用 VEIL_STOPS_DEFAULT;
    ang:   渐变角(1/60000 度), 缺省 GRADIENT_ANGLE_DOWN。
    调用顺序: 先 add_picture 铺底图, 再 add_veil, 最后放文字层。
    """
    if stops is None:
        stops = VEIL_STOPS_DEFAULT
    if ang is None:
        ang = GRADIENT_ANGLE_DOWN
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.line.fill.background()
    shape.shadow.inherit = False
    gs_xml = "".join(
        '<a:gs pos="%d"><a:srgbClr val="%s">%s</a:srgbClr></a:gs>' % (
            int(round(pos * 1000)), color,
            "" if opacity >= 100 else '<a:alpha val="%d"/>' % int(round(opacity * 1000)))
        for color, pos, opacity in stops)
    grad = parse_xml(
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' rotWithShape="1"><a:gsLst>%s</a:gsLst>'
        '<a:lin ang="%d" scaled="0"/></a:gradFill>' % (gs_xml, int(ang)))
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill",
                "a:blipFill", "a:pattFill", "a:grpFill"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    spPr.insert_element_before(grad, "a:ln", "a:effectLst", "a:effectDag",
                               "a:scene3d", "a:sp3d", "a:extLst")
    return shape

def add_header(slide, title, chapter=None, section=None):
    """页眉三件套: 绿竖条 + 29pt 粗白页题 (+ 右上面包屑「章  ·  节」)。

    chapter/section 缺省则只放竖条+页题; 破格页(S1/S2/S41 体例)别调本函数;
    章封面也别用它——用加高竖条(HEADER_BAR_Y_CHAPTER/H_CHAPTER)+章名+进度语变体。
    页题框定高 TITLE_H 且垂直居中(原片 spAutoFit 单行页 y/cy 会浮动,
    这里锚死恒定的垂直中心 TITLE_CY_CENTER, 视觉等价)。
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, HEADER_BAR_X, HEADER_BAR_Y, HEADER_BAR_W, HEADER_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()
    bar.shadow.inherit = False

    title_box = slide.shapes.add_textbox(TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = title
    set_font(r, SZ_TITLE, TEXT, bold=True)

    bread = None
    if chapter or section:
        bread = slide.shapes.add_textbox(BREAD_X, BREAD_Y, BREAD_W, BREAD_H)
        btf = bread.text_frame
        btf.word_wrap = True
        btf.auto_size = MSO_AUTO_SIZE.NONE
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        if chapter:
            r = p.add_run()
            r.text = chapter
            set_font(r, SZ_BREAD_CH, BRAND, bold=True)
        if chapter and section:
            r = p.add_run()
            r.text = BREAD_SEP
            set_font(r, SZ_BREAD_SEC, TEXT_DIM)
        if section:
            r = p.add_run()
            r.text = section
            set_font(r, SZ_BREAD_SEC, TEXT_DIM)
    return bar, title_box, bread

def add_conclusion_bar(slide, text):
    """底部结论条: 通栏 roundRect(PANEL 底 + 品牌绿 STROKE_BRAND 边), 收口一页的论点。

    文案含「：」时: 引导语(含冒号)16pt 绿粗 + 结论 15pt 白(S38 体例);
    不含则整句 16pt 绿粗(S8 裸排结论体例)。
    几何默认 S38 实测档(CONCL_*), y/高度随页内容可浮动——拿返回的 shape 自己改。
    原片 S38 边框 1.0pt 属描边漂移例外, v1.1 一律 1.2pt。
    """
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, CONCL_X, CONCL_Y, CONCL_W, CONCL_H)
    bar.adjustments[0] = 0.05          # adj=5000, 与金句框同档圆角
    bar.fill.solid()
    bar.fill.fore_color.rgb = PANEL
    bar.line.color.rgb = BRAND
    bar.line.width = STROKE_BRAND
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = CONCL_TEXT_INSET
    tf.margin_right = CONCL_TEXT_INSET
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if "：" in text:
        lead, rest = text.split("：", 1)
        r = p.add_run()
        r.text = lead + "："
        set_font(r, SZ_AGENDA, BRAND, bold=True)
        r = p.add_run()
        r.text = rest
        set_font(r, SZ_BODY, TEXT)
    else:
        r = p.add_run()
        r.text = text
        set_font(r, SZ_AGENDA, BRAND, bold=True)
    return bar

def _installed_font_families():
    """枚举 Windows 已安装字体族名(GDI EnumFontFamiliesExW, 含中文本地化名)。"""
    import ctypes
    from ctypes import wintypes

    gdi32 = ctypes.WinDLL("gdi32")
    user32 = ctypes.WinDLL("user32")

    class LOGFONTW(ctypes.Structure):
        _fields_ = [
            ("lfHeight", wintypes.LONG), ("lfWidth", wintypes.LONG),
            ("lfEscapement", wintypes.LONG), ("lfOrientation", wintypes.LONG),
            ("lfWeight", wintypes.LONG), ("lfItalic", wintypes.BYTE),
            ("lfUnderline", wintypes.BYTE), ("lfStrikeOut", wintypes.BYTE),
            ("lfCharSet", wintypes.BYTE), ("lfOutPrecision", wintypes.BYTE),
            ("lfClipPrecision", wintypes.BYTE), ("lfQuality", wintypes.BYTE),
            ("lfPitchAndFamily", wintypes.BYTE), ("lfFaceName", ctypes.c_wchar * 32),
        ]

    families = set()
    PROC = ctypes.WINFUNCTYPE(ctypes.c_int, ctypes.POINTER(LOGFONTW),
                              ctypes.c_void_p, wintypes.DWORD, wintypes.LPARAM)

    def _cb(lf, _tm, _font_type, _lparam):
        families.add(lf.contents.lfFaceName)
        return 1

    hdc = user32.GetDC(None)
    try:
        lf = LOGFONTW()
        lf.lfCharSet = 1               # DEFAULT_CHARSET: 枚举全部字符集
        gdi32.EnumFontFamiliesExW(hdc, ctypes.byref(lf), PROC(_cb), 0, 0)
    finally:
        user32.ReleaseDC(None, hdc)
    return families

def assert_fonts_installed(*fonts):
    """生成 deck 前跑一次: 枚举系统字体, 缺任一显式 raise, 绝不静默换字。

    缺省只查 FONT(思源黑体); 引入 mono 等新字体时把字体名一并传进来。
    PowerPoint 缺字时会无声替换字形, COM 渲染自检肉眼难辨——必须在生成侧堵死。
    """
    if not fonts:
        fonts = (FONT,)
    installed = _installed_font_families()
    missing = [f for f in fonts if f not in installed]
    if missing:
        raise RuntimeError(
            "渲染机缺字体: %s。请安装后重跑; 本库不做「微软雅黑」等静默兜底"
            "(typography 分册 Don't#4)。" % "、".join(missing))
